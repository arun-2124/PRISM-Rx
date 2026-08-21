import os
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------
# SETUP DIRECTORIES & CONSTANTS
# ---------------------------------------------------------
OUTPUT_DIR = r"d:\PRISM-Rx"
SCRATCH_DIR = os.path.join(OUTPUT_DIR, "scratch")
os.makedirs(SCRATCH_DIR, exist_ok=True)
PPTX_PATH = os.path.join(OUTPUT_DIR, "PRISM_Rx_Pitch_Deck.pptx")

# 16:9 Widescreen Dimensions
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Color Palette (Dark Mode Biotech AI Theme)
COLOR_BG = RGBColor(11, 15, 25)          # #0B0F19 (Dark Navy Canvas)
COLOR_CARD = RGBColor(19, 27, 46)        # #131B2E (Dark Slate Card)
COLOR_CARD_BORDER = RGBColor(30, 41, 59) # #1E293B (Card Border)
COLOR_TEXT_PRIMARY = RGBColor(248, 250, 252) # #F8FAFC (Off-White)
COLOR_TEXT_MUTED = RGBColor(148, 163, 184)   # #94A3B8 (Muted Gray)
COLOR_TEXT_SECONDARY = RGBColor(203, 213, 225) # #CBD5E1 (Light Slate)

COLOR_CYAN = RGBColor(0, 242, 254)       # #00F2FE (Primary Accent)
COLOR_VIOLET = RGBColor(139, 92, 246)    # #8B5CF6 (AI / Intelligence Accent)
COLOR_EMERALD = RGBColor(16, 185, 129)   # #10B981 (Success / High Score)
COLOR_ROSE = RGBColor(239, 68, 68)       # #EF4444 (Warning / Penalty Accent)

FONT_TITLE = "Trebuchet MS"
FONT_BODY = "Calibri"

# ---------------------------------------------------------
# MATPLOTLIB CHART GENERATION
# ---------------------------------------------------------
def generate_charts():
    plt.style.use('dark_background')
    
    # 1. Donut Chart: Node Distribution (1.31M Nodes)
    fig, ax = plt.subplots(figsize=(6, 4.5), facecolor='#131B2E')
    ax.set_facecolor('#131B2E')
    
    labels = ['Evidence\n(872.6K)', 'Clinical Trials\n(289.9K)', 'Targets\n(78.7K)', 'Diseases\n(47.1K)', 'Drugs\n(22.4K)']
    sizes = [872619, 289955, 78691, 47080, 22407]
    colors = ['#00F2FE', '#8B5CF6', '#10B981', '#F59E0B', '#EC4899']
    
    wedges, texts, autotexts = ax.pie(
        sizes, labels=labels, autopct='%1.1f%%', startangle=140,
        colors=colors, pctdistance=0.75,
        textprops=dict(color='#F8FAFC', fontsize=9, fontweight='bold')
    )
    for autotext in autotexts:
        autotext.set_color('#0B0F19')
        autotext.set_fontsize(9)
        
    centre_circle = plt.Circle((0,0), 0.55, fc='#131B2E')
    ax.add_artist(centre_circle)
    ax.set_title('Knowledge Graph Node Distribution (1.31M Total)', color='#00F2FE', fontsize=12, pad=12, fontweight='bold')
    
    chart_path1 = os.path.join(SCRATCH_DIR, "node_distribution.png")
    plt.tight_layout()
    plt.savefig(chart_path1, dpi=300, bbox_inches='tight', facecolor=fig.get_facecolor(), edgecolor='none')
    plt.close()

    # 2. Horizontal Bar Chart: Path Conversion & Candidates
    fig, ax = plt.subplots(figsize=(6.5, 4.5), facecolor='#131B2E')
    ax.set_facecolor('#131B2E')
    
    metrics = ['Total Graph Paths', 'Unindicated Paths', 'Unique Pairs']
    values = [2976634, 2787952, 819696]
    bar_colors = ['#8B5CF6', '#00F2FE', '#10B981']
    
    y_pos = np.arange(len(metrics))
    bars = ax.barh(y_pos, values, color=bar_colors, height=0.55)
    
    ax.set_yticks(y_pos)
    ax.set_yticklabels(metrics, color='#F8FAFC', fontsize=10, fontweight='bold')
    ax.invert_yaxis()  # top-down
    ax.set_xlabel('Count (Millions)', color='#94A3B8', fontsize=10)
    ax.set_title('Graph Traversals to Evaluated Candidates', color='#00F2FE', fontsize=12, pad=12, fontweight='bold')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#334155')
    ax.spines['left'].set_color('#334155')
    ax.tick_params(colors='#94A3B8')
    
    for bar in bars:
        width = bar.get_width()
        ax.text(width + 50000, bar.get_y() + bar.get_height()/2, f'{width:,}',
                va='center', ha='left', color='#F8FAFC', fontsize=9, fontweight='bold')
                
    chart_path2 = os.path.join(SCRATCH_DIR, "path_conversion.png")
    plt.tight_layout()
    plt.savefig(chart_path2, dpi=300, bbox_inches='tight', facecolor=fig.get_facecolor(), edgecolor='none')
    plt.close()

    return chart_path1, chart_path2

# ---------------------------------------------------------
# PPTX BUILDER HELPERS
# ---------------------------------------------------------
def create_blank_slide(prs):
    blank_slide_layout = prs.slide_layouts[6] # blank
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background rectangle
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG
    bg.line.fill.background()
    return slide

def add_header(slide, title_text, subtitle_text):
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = subtitle_text.upper()
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.name = FONT_TITLE
    p.font.color.rgb = COLOR_CYAN
    
    p2 = tf.add_paragraph()
    p2.text = title_text
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.name = FONT_TITLE
    p2.font.color.rgb = COLOR_TEXT_PRIMARY

def add_card(slide, left, top, width, height, bg_color=COLOR_CARD, border_color=COLOR_CARD_BORDER):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
    else:
        card.line.fill.background()
    return card

def add_kpi_card(slide, left, top, width, height, value, label, subtext="", accent_color=COLOR_CYAN):
    add_card(slide, left, top, width, height)
    
    # Top color indicator bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = FONT_TITLE
    p.font.color.rgb = accent_color
    
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.font.name = FONT_BODY
    p2.font.color.rgb = COLOR_TEXT_PRIMARY
    
    if subtext:
        p3 = tf.add_paragraph()
        p3.text = subtext
        p3.font.size = Pt(10)
        p3.font.name = FONT_BODY
        p3.font.color.rgb = COLOR_TEXT_MUTED

# ---------------------------------------------------------
# SLIDE BUILDERS (16 SLIDES)
# ---------------------------------------------------------

def build_slide_1(prs): # TITLE
    slide = create_blank_slide(prs)
    
    # Hero Center Container
    add_card(slide, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1), bg_color=COLOR_CARD, border_color=COLOR_CYAN)
    
    # Category Pill
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.666), Inches(1.6), Inches(4.0), Inches(0.4))
    pill.fill.solid()
    pill.fill.fore_color.rgb = COLOR_CARD_BORDER
    pill.line.color.rgb = COLOR_CYAN
    tf_p = pill.text_frame
    tf_p.text = "BIOTECH RESEARCH INTELLIGENCE"
    p_p = tf_p.paragraphs[0]
    p_p.alignment = PP_ALIGN.CENTER
    p_p.font.size = Pt(10)
    p_p.font.bold = True
    p_p.font.color.rgb = COLOR_CYAN
    
    # Title & Subtitle Text
    txBox = slide.shapes.add_textbox(Inches(2.0), Inches(2.2), Inches(9.333), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "PRISM-Rx"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.size = Pt(48)
    p1.font.bold = True
    p1.font.name = FONT_TITLE
    p1.font.color.rgb = COLOR_TEXT_PRIMARY
    
    p2 = tf.add_paragraph()
    p2.text = "Real-Time Biotech Arbitrage Engine for Drug Repurposing Signals"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.name = FONT_BODY
    p2.font.color.rgb = COLOR_CYAN
    
    p3 = tf.add_paragraph()
    p3.text = "Automated Multi-Hop Knowledge Graph Traversals & 100-Point Multi-Factor Evidence Scoring"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(12)
    p3.font.name = FONT_BODY
    p3.font.color.rgb = COLOR_TEXT_MUTED
    
    # 3 Badges at Bottom of Card
    badge_w = Inches(2.8)
    badge_h = Inches(0.9)
    top_b = Inches(4.9)
    
    add_kpi_card(slide, Inches(2.2), top_b, badge_w, badge_h, "1.31M+ Nodes", "Knowledge Graph", "Open Targets, ChEMBL, PMC", COLOR_CYAN)
    add_kpi_card(slide, Inches(5.266), top_b, badge_w, badge_h, "819K+ Candidates", "Evaluated Pairs", "Unique (Drug, Disease)", COLOR_VIOLET)
    add_kpi_card(slide, Inches(8.333), top_b, badge_w, badge_h, "FastAPI + React", "Production Stack", "Real-Time Interactive UI", COLOR_EMERALD)

def build_slide_2(prs): # THE PROBLEM
    slide = create_blank_slide(prs)
    add_header(slide, "Drug Discovery Bottlenecks & Siloed Evidence", "The Problem")
    
    # 3 KPI Cards Top Row
    card_w = Inches(3.64)
    top_pos = Inches(1.5)
    add_kpi_card(slide, Inches(0.8), top_pos, card_w, Inches(1.3), "$2.6 Billion", "Average Cost Per Drug", "De novo drug discovery capital expenditure", COLOR_ROSE)
    add_kpi_card(slide, Inches(4.84), top_pos, card_w, Inches(1.3), "10–15 Years", "Average Time to Market", "Extensive clinical trial & approval cycles", COLOR_ROSE)
    add_kpi_card(slide, Inches(8.88), top_pos, card_w, Inches(1.3), "90% Failure Rate", "Clinical Pipeline Attrition", "High failure due to safety or efficacy gaps", COLOR_ROSE)
    
    # 2 Main Problem Cards Bottom
    bot_top = Inches(3.1)
    bot_h = Inches(3.8)
    w_main = Inches(5.66)
    
    # Card 1
    add_card(slide, Inches(0.8), bot_top, w_main, bot_h)
    tx1 = slide.shapes.add_textbox(Inches(1.1), bot_top + Inches(0.3), w_main - Inches(0.6), bot_h - Inches(0.6))
    tf1 = tx1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "CRITICAL PROBLEM 1"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_ROSE
    
    p2 = tf1.add_paragraph()
    p2.text = "Fragmented Biomedical Data Silos"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT_PRIMARY
    
    points1 = [
        "Biomedical data is split across Open Targets, ChEMBL, PubMed, Europe PMC, UniProt, and ClinicalTrials.gov.",
        "Researchers manually consult isolated tools, missing subtle cross-database mechanistic links.",
        "No existing open system unifies target associations, chemical structures, and clinical reports into a single multi-hop graph."
    ]
    for pt in points1:
        p_pt = tf1.add_paragraph()
        p_pt.text = f"•  {pt}"
        p_pt.font.size = Pt(12)
        p_pt.font.color.rgb = COLOR_TEXT_SECONDARY
        
    # Card 2
    add_card(slide, Inches(6.866), bot_top, w_main, bot_h)
    tx2 = slide.shapes.add_textbox(Inches(7.166), bot_top + Inches(0.3), w_main - Inches(0.6), bot_h - Inches(0.6))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    
    p = tf2.paragraphs[0]
    p.text = "CRITICAL PROBLEM 2"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_ROSE
    
    p2 = tf2.add_paragraph()
    p2.text = "Overlooked Repurposing Signals"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT_PRIMARY
    
    points2 = [
        "Over 2.78 Million unindicated Drug-Target-Disease pathways exist in public data but remain unanalyzed.",
        "Traditional drug discovery ignores approved, safe compounds that already possess Phase 1 safety profiles.",
        "Manual research cannot evaluate 819,000+ potential drug-disease candidate pairs against multi-factor evidence."
    ]
    for pt in points2:
        p_pt = tf2.add_paragraph()
        p_pt.text = f"•  {pt}"
        p_pt.font.size = Pt(12)
        p_pt.font.color.rgb = COLOR_TEXT_SECONDARY

def build_slide_3(prs): # WHY THE PROBLEM MATTERS
    slide = create_blank_slide(prs)
    add_header(slide, "Clinical Urgency & Economic Consequences", "Why The Problem Matters")
    
    card_w = Inches(3.64)
    card_h = Inches(5.2)
    top_pos = Inches(1.5)
    
    cols = [
        ("Clinical Urgency", COLOR_CYAN, [
            "Patients with rare or aggressive cancers (e.g. Acute Lymphoblastic Leukemia) cannot wait 10+ years for de novo drug discovery.",
            "Repurposing safe, approved drugs delivers immediate therapeutic candidates to clinical trials.",
            "Significantly reduces early-stage toxicology failures by leveraging established human safety data."
        ]),
        ("R&D Economics", COLOR_VIOLET, [
            "Pharma R&D productivity has steadily declined despite exponential growth in biomedical literature.",
            "Repurposing cuts R&D costs by up to 60% and development timelines by 5–7 years.",
            "Maximizes return on past capital investment by discovering new patentable disease indications for existing assets."
        ]),
        ("Scale of Data", COLOR_EMERALD, [
            "Human researchers can read ~500 papers per year; over 1.5 million life science papers are published annually.",
            "Manual review creates severe blind spots, overlooking novel multi-target drug mechanisms.",
            "Requires automated graph intelligence to continuously evaluate evidence across millions of biological entities."
        ])
    ]
    
    lefts = [Inches(0.8), Inches(4.84), Inches(8.88)]
    for idx, (title, color, bullets) in enumerate(cols):
        add_card(slide, lefts[idx], top_pos, card_w, card_h)
        
        # Color top line
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, lefts[idx], top_pos, card_w, Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        tx = slide.shapes.add_textbox(lefts[idx] + Inches(0.3), top_pos + Inches(0.3), card_w - Inches(0.6), card_h - Inches(0.6))
        tf = tx.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_PRIMARY
        
        for b in bullets:
            pb = tf.add_paragraph()
            pb.text = f"•  {b}"
            pb.font.size = Pt(12)
            pb.font.color.rgb = COLOR_TEXT_SECONDARY

def build_slide_4(prs): # EXISTING APPROACH / GAP
    slide = create_blank_slide(prs)
    add_header(slide, "Current Approaches vs. Systemic Limitations", "Existing Approach & Gaps")
    
    card_w = Inches(3.64)
    card_h = Inches(5.2)
    top_pos = Inches(1.5)
    lefts = [Inches(0.8), Inches(4.84), Inches(8.88)]
    
    approaches = [
        ("Manual Literature Review", "TRADITIONAL APPROACH", COLOR_ROSE, [
            "Status: Researchers perform manual keyword searches across PubMed & Google Scholar.",
            "Limitation 1: High human error and extreme cognitive overload.",
            "Limitation 2: Cannot connect indirect 2-hop or 3-hop biological paths across disparate papers.",
            "Result: Overlooks 95%+ of novel repurposing opportunities."
        ]),
        ("Single-Source Database Queries", "ISOLATED TOOLS", COLOR_ROSE, [
            "Status: Querying standalone portals like Open Targets or ChEMBL.",
            "Limitation 1: Datasets remain un-linked; target binding does not automatically map to clinical trials.",
            "Limitation 2: No unified evidence scoring across literature, trials, and target confidence.",
            "Result: Fragmented insights requiring tedious manual synthesis."
        ]),
        ("Naïve Unweighted Graphs", "BASIC GRAPH SEARCH", COLOR_ROSE, [
            "Status: Running standard relational SQL or graph queries without scoring engines.",
            "Limitation 1: Severe path duplication (2.97M paths for 819K pairs) floods users with redundant rows.",
            "Limitation 2: No safety penalty checks; includes withdrawn or black-box toxic drugs.",
            "Result: Noise overwhelms actionable research signals."
        ])
    ]
    
    for idx, (title, badge, color, bullets) in enumerate(approaches):
        add_card(slide, lefts[idx], top_pos, card_w, card_h)
        
        tx = slide.shapes.add_textbox(lefts[idx] + Inches(0.3), top_pos + Inches(0.3), card_w - Inches(0.6), card_h - Inches(0.6))
        tf = tx.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = badge
        p0.font.size = Pt(10)
        p0.font.bold = True
        p0.font.color.rgb = color
        
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_PRIMARY
        
        for b in bullets:
            pb = tf.add_paragraph()
            pb.text = b
            pb.font.size = Pt(11)
            pb.font.color.rgb = COLOR_TEXT_SECONDARY

def build_slide_5(prs): # OUR SOLUTION
    slide = create_blank_slide(prs)
    add_header(slide, "PRISM-Rx: Real-Time Biotech Arbitrage Engine", "Our Solution")
    
    # Hero Central Card
    add_card(slide, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.2), bg_color=COLOR_CARD, border_color=COLOR_CYAN)
    
    tx = slide.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(11.133), Inches(4.6))
    tf = tx.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "A Unified Research Intelligence Layer for Drug Repurposing"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    
    # 4 Value Pillars in 2x2 Grid
    pillars = [
        ("1.31M+ Node Knowledge Graph", "Integrates 5 public databases into a single relational SQLite layer (medbase.db, 545 MB) connecting Drugs, Targets, Diseases, Clinical Trials, and Literature."),
        ("Duplicate Candidate Path Collapsing", "Engine V2/V3 collapses 2.97M raw graph paths into 819,696 unique (Drug, Disease) pairs with multi-target log bonus calculations."),
        ("100-Point Multi-Factor Evidence Model", "Evaluates candidate pairs across 7 positive evidence dimensions (Target score, Drug-Target confidence, Clinical phase, Literature quality, Diversity, Novelty)."),
        ("Dynamic Safety & Contradiction Penalties", "Automated deductions (-40 pts for withdrawn drugs, -25 pts for black-box warnings, -30 pts for biological mechanism contradictions).")
    ]
    
    for idx, (title, desc) in enumerate(pillars):
        p_t = tf.add_paragraph()
        p_t.text = f"\n✦  {title}"
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT_PRIMARY
        
        p_d = tf.add_paragraph()
        p_d.text = f"    {desc}"
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_TEXT_SECONDARY

def build_slide_6(prs): # HOW IT WORKS
    slide = create_blank_slide(prs)
    add_header(slide, "End-to-End Data Ingestion & Signal Synthesis", "How It Works")
    
    card_w = Inches(2.7)
    card_h = Inches(5.2)
    top_pos = Inches(1.5)
    
    steps = [
        ("STEP 1", "Multi-Source Ingestion", COLOR_CYAN, [
            "HTTPS Parquet downloads from Open Targets 26.06.",
            "SQLite dump extraction from ChEMBL 37.",
            "REST API fetches from Europe PMC, UniProt, and ClinicalTrials.gov v2."
        ]),
        ("STEP 2", "Entity Normalization", COLOR_VIOLET, [
            "Authoritative cross-reference mapping (no fuzzy matching).",
            "Drugs: ChEMBL IDs, Open Targets IDs.",
            "Targets: UniProt IDs, Ensembl IDs.",
            "Diseases: EFO, MONDO ontologies."
        ]),
        ("STEP 3", "Knowledge Graph Storage", COLOR_EMERALD, [
            "SQLite database (medbase.db, 545 MB).",
            "1.31M+ Nodes & 1.08M+ Edges.",
            "10 indexed core relational entity tables with strict constraints."
        ]),
        ("STEP 4", "Signal Scoring & UI", COLOR_CYAN, [
            "Candidate collapsing into 819.6K unique pairs.",
            "100-Point multi-factor evidence scoring & safety penalties.",
            "FastAPI REST API + React Vite interactive graph dashboard."
        ])
    ]
    
    lefts = [Inches(0.8), Inches(3.8), Inches(6.8), Inches(9.8)]
    for idx, (step, title, color, bullets) in enumerate(steps):
        add_card(slide, lefts[idx], top_pos, card_w, card_h)
        
        tx = slide.shapes.add_textbox(lefts[idx] + Inches(0.2), top_pos + Inches(0.2), card_w - Inches(0.4), card_h - Inches(0.4))
        tf = tx.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = step
        p0.font.size = Pt(10)
        p0.font.bold = True
        p0.font.color.rgb = color
        
        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_PRIMARY
        
        for b in bullets:
            pb = tf.add_paragraph()
            pb.text = f"• {b}"
            pb.font.size = Pt(11)
            pb.font.color.rgb = COLOR_TEXT_SECONDARY

def build_slide_7(prs): # KEY FEATURES
    slide = create_blank_slide(prs)
    add_header(slide, "Platform Capabilities & Scientific Tools", "Key Features")
    
    card_w = Inches(3.64)
    card_h = Inches(2.4)
    
    features = [
        ("Candidate Path Collapsing", "Collapses 2.97M redundant graph paths into unique (Drug, Disease) candidate pairs with multi-target log bonuses.", COLOR_CYAN),
        ("100-Point Evidence Scoring", "Transparent multi-factor scoring combining target association, action confidence, clinical phase, literature, diversity, and novelty.", COLOR_VIOLET),
        ("Safety & Toxicity Deductions", "Automated safety penalty engine deducting up to -40 pts for black-box warnings, drug withdrawals, and mechanism contradictions.", COLOR_ROSE),
        ("Interactive 2-Hop Force Graph", "Custom React HTML5/SVG canvas graph rendering Drug -> Target -> Disease topology with drag, zoom, and node inspection.", COLOR_EMERALD),
        ("RAG Copilot Grounding", "Natural language chat assistant grounded strictly in 32+ provenanced evidence records from medbase.db with zero hallucination.", COLOR_CYAN),
        ("Multi-Format Export", "One-click export of prioritized candidate reports, score breakdowns, and biological paths to formatted CSV or JSON payloads.", COLOR_VIOLET)
    ]
    
    grid_coords = [
        (Inches(0.8), Inches(1.5)), (Inches(4.84), Inches(1.5)), (Inches(8.88), Inches(1.5)),
        (Inches(0.8), Inches(4.2)), (Inches(4.84), Inches(4.2)), (Inches(8.88), Inches(4.2))
    ]
    
    for idx, (title, desc, color) in enumerate(features):
        left, top = grid_coords[idx]
        add_card(slide, left, top, card_w, card_h)
        
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, card_w, Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        tx = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), card_w - Inches(0.4), card_h - Inches(0.4))
        tf = tx.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_PRIMARY
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_SECONDARY

def build_slide_8(prs): # TECHNICAL ARCHITECTURE
    slide = create_blank_slide(prs)
    add_header(slide, "Full-Stack System Architecture Specification", "System Architecture")
    
    card_w = Inches(11.733)
    card_h = Inches(1.15)
    left_pos = Inches(0.8)
    
    layers = [
        ("LAYER 1: REACT + VITE FRONTEND (PORT 3000)", COLOR_CYAN,
         "Dashboard Stat Cards | Signal Explorer Filter Sidebar | Interactive 2-Hop Canvas Graph | RAG Copilot UI | Methodology & Specs"),
        ("LAYER 2: FASTAPI BACKEND SERVICE (PORT 8000)", COLOR_VIOLET,
         "12 REST Endpoints: /api/signals, /api/signals/{id}, /api/graph/{id}, /api/clinical-trials/{id}, /api/evidence/{id}, /api/export"),
        ("LAYER 3: SIGNAL INTELLIGENCE ENGINE V3 & GRAPH TRAVERSAL ENGINE", COLOR_EMERALD,
         "SignalEngineV3 (Candidate Collapsing & 100-Pt Scoring) | GraphTraversalEngine (1-hop/2-hop Neighborhoods & Provenance Tracking)"),
        ("LAYER 4: UNIFIED SQLITE DATABASE (medbase.db)", COLOR_CYAN,
         "545.35 MB Relational File | 2,002,249 Total Records | 16 Indexed Tables (drugs, diseases, targets, clinical_reports, evidence, drug_warnings)")
    ]
    
    tops = [Inches(1.5), Inches(2.85), Inches(4.2), Inches(5.55)]
    
    for idx, (title, color, text) in enumerate(layers):
        add_card(slide, left_pos, tops[idx], card_w, card_h)
        
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, tops[idx], Inches(0.1), card_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        tx = slide.shapes.add_textbox(left_pos + Inches(0.3), tops[idx] + Inches(0.15), card_w - Inches(0.5), card_h - Inches(0.3))
        tf = tx.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = text
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_PRIMARY

def build_slide_9(prs): # AI / INTELLIGENCE
    slide = create_blank_slide(prs)
    add_header(slide, "100-Point Multi-Factor Evidence Scoring Formula", "AI / Scoring Intelligence")
    
    # Formula Box Top
    add_card(slide, Inches(0.8), Inches(1.5), Inches(11.733), Inches(1.1), bg_color=COLOR_CARD, border_color=COLOR_CYAN)
    tx_f = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.333), Inches(0.9))
    tf_f = tx_f.text_frame
    tf_f.word_wrap = True
    
    p = tf_f.paragraphs[0]
    p.text = "RESEARCH PRIORITY SCORE FORMULA"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    
    p2 = tf_f.add_paragraph()
    p2.text = "Score = Clamp( S_TD + S_DT + S_Clin + S_Lit + F_Div + B_Target + S_Nov - P_Safety - P_Contra, 0, 100 )"
    p2.font.size = Pt(15)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT_PRIMARY
    
    # 2 Column Cards Below
    w_col = Inches(5.66)
    top_c = Inches(2.8)
    h_c = Inches(4.2)
    
    # Positive Dimensions
    add_card(slide, Inches(0.8), top_c, w_col, h_c)
    tx_p = slide.shapes.add_textbox(Inches(1.0), top_c + Inches(0.2), w_col - Inches(0.4), h_c - Inches(0.4))
    tf_p = tx_p.text_frame
    tf_p.word_wrap = True
    
    p = tf_p.paragraphs[0]
    p.text = "POSITIVE EVIDENCE DIMENSIONS (+100 PTS MAX)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_EMERALD
    
    pos_dims = [
        ("Target-Disease Association (S_TD)", "Max 30 pts — Based on highest experimental score"),
        ("Drug-Target Action Confidence (S_DT)", "Max 15 pts — Inhibitor (0.9), Agonist (0.8), Modulator (0.7)"),
        ("Clinical Evidence & Phase (S_Clin)", "Max 15 pts — Phase 4 (1.0), Phase 3 (0.8), Phase 2 (0.6), Phase 1 (0.4)"),
        ("Literature & Evidence Tier (S_Lit)", "Max 10 pts — High Quality (1.0), Medium (0.7), Low (0.4)"),
        ("Source Diversity Factor (F_Div)", "Max 10 pts — Verified across Open Targets, Europe PMC, Trials, UniProt"),
        ("Multi-Target Support Bonus (B_Target)", "Max 10 pts — Capped log bonus for multi-target hits: 5.0 * log2(N)"),
        ("Under-Investigated Novelty (S_Nov)", "Max 10 pts — Non-indicated candidate status reward")
    ]
    for name, desc in pos_dims:
        pb = tf_p.add_paragraph()
        pb.text = f"• {name}: {desc}"
        pb.font.size = Pt(10)
        pb.font.color.rgb = COLOR_TEXT_SECONDARY
        
    # Negative Penalties
    add_card(slide, Inches(6.866), top_c, w_col, h_c)
    tx_n = slide.shapes.add_textbox(Inches(7.066), top_c + Inches(0.2), w_col - Inches(0.4), h_c - Inches(0.4))
    tf_n = tx_n.text_frame
    tf_n.word_wrap = True
    
    p = tf_n.paragraphs[0]
    p.text = "SAFETY & CONTRADICTION PENALTIES"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_ROSE
    
    neg_dims = [
        ("Safety & Black-Box Warning (P_Safety)", "Deduction -10 to -40 pts — Withdrawn status (-40 pts), Black-box warning (-25 pts), General warning (-10 pts)."),
        ("Biological Contradiction (P_Contra)", "Deduction -15 to -30 pts — Opposing direction on disease trait or conflicting mechanism of action."),
        ("Categorization Tiers", "• STRONG_RESEARCH_SIGNAL: Score >= 70.0 (0 penalties)\n• MODERATE_RESEARCH_SIGNAL: 40.0 <= Score < 70.0\n• WEAK_RESEARCH_SIGNAL: 20.0 <= Score < 40.0\n• CONTRADICTED: Penalty >= 25.0 or Contradiction > 0")
    ]
    for name, desc in neg_dims:
        pb = tf_n.add_paragraph()
        pb.text = f"• {name}:\n  {desc}"
        pb.font.size = Pt(10)
        pb.font.color.rgb = COLOR_TEXT_SECONDARY

def build_slide_10(prs, chart1, chart2): # DATA / ANALYTICS
    slide = create_blank_slide(prs)
    add_header(slide, "Knowledge Graph Analytics & Top Candidate Signal", "Data & Signals")
    
    # Left: Donut Chart
    slide.shapes.add_picture(chart1, Inches(0.8), Inches(1.5), width=Inches(5.6))
    
    # Right Top: Top Signal Highlight Card
    add_card(slide, Inches(6.6), Inches(1.5), Inches(5.933), Inches(2.6), bg_color=COLOR_CARD, border_color=COLOR_EMERALD)
    tx_s = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.533), Inches(2.4))
    tf_s = tx_s.text_frame
    tf_s.word_wrap = True
    
    p = tf_s.paragraphs[0]
    p.text = "TOP GROUNDED CANDIDATE SIGNAL CASE STUDY"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLOR_EMERALD
    
    p2 = tf_s.add_paragraph()
    p2.text = "Tg100-801 -> Acute Lymphoblastic Leukemia"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT_PRIMARY
    
    p3 = tf_s.add_paragraph()
    p3.text = "• PRISM Research Priority Score: 82.0 / 100 (STRONG_RESEARCH_SIGNAL)\n• Arbitrage Score: 9.3 / 10.0 | Convergence Score: 81.3 / 100\n• Target: FGR proto-oncogene (SRC inhibitor | Action Conf: 0.9)\n• Provenance: 3 Independent Sources | 32 Provenanced Records"
    p3.font.size = Pt(11)
    p3.font.color.rgb = COLOR_TEXT_SECONDARY
    
    # Right Bottom: Path Conversion Chart
    slide.shapes.add_picture(chart2, Inches(6.6), Inches(4.3), width=Inches(5.933))

def build_slide_11(prs): # IMPACT
    slide = create_blank_slide(prs)
    add_header(slide, "Empirical Platform Benchmarks & Knowledge Scale", "Impact & Performance")
    
    w_card = Inches(5.66)
    h_card = Inches(2.4)
    
    kpis = [
        ("1,310,752", "Aggregated Knowledge Graph Nodes", "Unified across Open Targets, ChEMBL, PubMed, Europe PMC, UniProt, and ClinicalTrials.gov v2.", COLOR_CYAN, Inches(0.8), Inches(1.5)),
        ("819,696", "Unique Repurposing Candidates Evaluated", "Derived from collapsing 2.97M raw multi-hop graph paths into distinct (Drug, Disease) pairs.", COLOR_VIOLET, Inches(6.866), Inches(1.5)),
        ("93.66%", "Novel Unindicated Pathways Uncovered", "2.78M unindicated candidate paths prioritized out of 2.97M total traversals.", COLOR_EMERALD, Inches(0.8), Inches(4.2)),
        ("3.34 Seconds", "End-to-End Analytics & Traversal Speed", "High-performance indexed SQLite query execution across 2.0M+ database records.", COLOR_CYAN, Inches(6.866), Inches(4.2))
    ]
    
    for val, title, desc, color, left, top in kpis:
        add_kpi_card(slide, left, top, w_card, h_card, val, title, desc, color)

def build_slide_12(prs): # IMPLEMENTATION / DEPLOYMENT
    slide = create_blank_slide(prs)
    add_header(slide, "System Development & Engineering Roadmap", "Implementation Phases")
    
    card_w = Inches(11.733)
    card_h = Inches(1.15)
    left_pos = Inches(0.8)
    
    phases = [
        ("PHASE 1-2: DATA INGESTION & ETL PIPELINE", COLOR_CYAN,
         "Ingested Open Targets 26.06 Parquet datasets, ChEMBL 37 SQLite database dump, Europe PMC REST API, UniProt API, ClinicalTrials.gov v2."),
        ("PHASE 3-4: ENTITY NORMALIZATION & UNIFIED DB", COLOR_VIOLET,
         "Built strict cross-reference normalization mapping (ChEMBL, UniProt, Ensembl, EFO/MONDO). Created unified SQLite DB (medbase.db, 545 MB)."),
        ("PHASE 5-6: SIGNAL INTELLIGENCE ENGINE V2/V3", COLOR_EMERALD,
         "Implemented duplicate path collapsing into 819.6K candidate pairs, 100-point multi-factor evidence scoring model, and safety penalty engine."),
        ("PHASE 7-8: FULL-STACK DASHBOARD & RAG COPILOT", COLOR_CYAN,
         "Deployed FastAPI REST service (12 endpoints), React + Vite single-page dashboard with interactive 2-hop graph, and RAG copilot integration.")
    ]
    
    tops = [Inches(1.5), Inches(2.85), Inches(4.2), Inches(5.55)]
    
    for idx, (title, color, text) in enumerate(phases):
        add_card(slide, left_pos, tops[idx], card_w, card_h)
        
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, tops[idx], Inches(0.1), card_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        tx = slide.shapes.add_textbox(left_pos + Inches(0.3), tops[idx] + Inches(0.15), card_w - Inches(0.5), card_h - Inches(0.3))
        tf = tx.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = text
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_PRIMARY

def build_slide_13(prs): # SCALABILITY / FUTURE
    slide = create_blank_slide(prs)
    add_header(slide, "Scalability & Next-Generation Platform Horizons", "Business & Scalability")
    
    card_w = Inches(5.66)
    card_h = Inches(2.4)
    
    items = [
        ("Real-Time Ingestion Webhooks", "Automating real-time streaming ingestion for daily PubMed literature preprints and ClinicalTrials.gov study status updates.", COLOR_CYAN, Inches(0.8), Inches(1.5)),
        ("LLM NLP Mechanism Extraction", "Integrating specialized biomedical NLP models to automatically extract directional protein-target mechanism actions from un-annotated text.", COLOR_VIOLET, Inches(6.866), Inches(1.5)),
        ("Multi-Omic Data Expansion", "Incorporating GTEx tissue eQTL expression quantitative trait loci, Single-cell RNA-seq, and ENCODE cCRE cis-regulatory elements.", COLOR_EMERALD, Inches(0.8), Inches(4.2)),
        ("Enterprise Biotech Integration", "Deploying containerized Docker microservices for institutional pharmaceutical R&D cloud deployment and API monetization.", COLOR_CYAN, Inches(6.866), Inches(4.2))
    ]
    
    for title, desc, color, left, top in items:
        add_card(slide, left, top, card_w, card_h)
        
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, card_w, Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        tx = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.2), card_w - Inches(0.6), card_h - Inches(0.4))
        tf = tx.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_PRIMARY
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_SECONDARY

def build_slide_14(prs): # COMPETITIVE ADVANTAGE
    slide = create_blank_slide(prs)
    add_header(slide, "System Differentiation vs. Existing Approaches", "Competitive Advantage")
    
    # Table shape
    rows, cols = 6, 5
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.2))
    table = table_shape.table
    
    # Column widths
    table.columns[0].width = Inches(2.933)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(2.2)
    table.columns[3].width = Inches(2.2)
    table.columns[4].width = Inches(2.2)
    
    headers = ["FEATURE / CAPABILITY", "MANUAL PUBMED", "STANDALONE DBS", "GENERIC LLMS", "PRISM-Rx PLATFORM"]
    data = [
        ["Multi-Hop Graph Traversals", "Manual (Impossible)", "No (1-hop only)", "No (Hallucinates)", "Automated (1.31M Nodes)"],
        ["Duplicate Candidate Collapsing", "None", "None", "None", "Collapses 2.97M -> 819.6K"],
        ["100-Pt Multi-Factor Scoring", "None", "Uniform / Static", "None", "7 Evidence + 2 Penalty Factors"],
        ["Safety & Contradiction Penalties", "Manual check", "Ignored", "Ignored", "Automated (-40 / -30 pts)"],
        ["Grounded Provenance & Audit", "High Human Error", "Partial", "Frequent Hallucination", "100% Provenanced DB Records"]
    ]
    
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_CARD_BORDER if col_idx < 4 else COLOR_CYAN
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_BG if col_idx == 4 else COLOR_TEXT_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
    for row_idx, row_data in enumerate(data):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            if col_idx == 4:
                cell.fill.fore_color.rgb = COLOR_CARD
            else:
                cell.fill.fore_color.rgb = COLOR_BG
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(10)
            p.font.bold = (col_idx == 4 or col_idx == 0)
            p.font.color.rgb = COLOR_CYAN if col_idx == 4 else COLOR_TEXT_SECONDARY
            if col_idx > 0:
                p.alignment = PP_ALIGN.CENTER

def build_slide_15(prs): # CONCLUSION
    slide = create_blank_slide(prs)
    add_header(slide, "Transforming Biomedical Data into Actionable Signals", "Conclusion")
    
    # 3 Summary Cards
    card_w = Inches(3.64)
    card_h = Inches(3.6)
    top_pos = Inches(1.5)
    lefts = [Inches(0.8), Inches(4.84), Inches(8.88)]
    
    cards_data = [
        ("Unified Evidence Layer", COLOR_CYAN, [
            "Replaces fragmented biomedical silos with a 1.31M+ node multi-hop Knowledge Graph.",
            "Connects Open Targets, ChEMBL, PubMed, UniProt, and ClinicalTrials.gov."
        ]),
        ("Explainable Intelligence", COLOR_VIOLET, [
            "100-Point multi-factor evidence scoring model delivers full mathematical transparency.",
            "Enforces strict safety and contradiction penalty deductions."
        ]),
        ("Accelerated Discovery", COLOR_EMERALD, [
            "Unlocks 819,696 unique repurposing candidate signals in under 3.5 seconds.",
            "Prioritizes de-risked safe compounds for rapid clinical translation."
        ])
    ]
    
    for idx, (title, color, bullets) in enumerate(cards_data):
        add_card(slide, lefts[idx], top_pos, card_w, card_h)
        
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, lefts[idx], top_pos, card_w, Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        tx = slide.shapes.add_textbox(lefts[idx] + Inches(0.3), top_pos + Inches(0.3), card_w - Inches(0.6), card_h - Inches(0.6))
        tf = tx.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_PRIMARY
        
        for b in bullets:
            pb = tf.add_paragraph()
            pb.text = f"•  {b}"
            pb.font.size = Pt(12)
            pb.font.color.rgb = COLOR_TEXT_SECONDARY
            
    # Closing Banner Bottom
    add_card(slide, Inches(0.8), Inches(5.4), Inches(11.733), Inches(1.3), bg_color=COLOR_CARD, border_color=COLOR_CYAN)
    tx_b = slide.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(11.333), Inches(1.1))
    tf_b = tx_b.text_frame
    tf_b.word_wrap = True
    
    p_b = tf_b.paragraphs[0]
    p_b.text = "PRISM-Rx transforms millions of fragmented biological data points into high-confidence, actionable drug repurposing signals."
    p_b.alignment = PP_ALIGN.CENTER
    p_b.font.size = Pt(18)
    p_b.font.bold = True
    p_b.font.color.rgb = COLOR_CYAN

def build_slide_16(prs): # THANK YOU
    slide = create_blank_slide(prs)
    
    # Center Hero Card
    add_card(slide, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1), bg_color=COLOR_CARD, border_color=COLOR_CYAN)
    
    txBox = slide.shapes.add_textbox(Inches(2.0), Inches(1.6), Inches(9.333), Inches(4.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "THANK YOU"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.name = FONT_TITLE
    p1.font.color.rgb = COLOR_CYAN
    
    p2 = tf.add_paragraph()
    p2.text = "PRISM-Rx: Real-Time Biotech Arbitrage Engine"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT_PRIMARY
    
    p3 = tf.add_paragraph()
    p3.text = "\nProject Codebase & Documentation: d:\\PRISM-Rx\nFastAPI REST Service & React Vite Dashboard Ready"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_TEXT_SECONDARY
    
    p4 = tf.add_paragraph()
    p4.text = "\nScientific Disclaimer: All outputs are Computational Research Hypotheses / Research Priorities.\nPRISM-Rx does NOT generate medical recommendations or predict clinical efficacy."
    p4.alignment = PP_ALIGN.CENTER
    p4.font.size = Pt(10)
    p4.font.color.rgb = COLOR_TEXT_MUTED

# ---------------------------------------------------------
# MAIN EXECUTION
# ---------------------------------------------------------
def main():
    print("Generating Matplotlib charts...")
    chart1, chart2 = generate_charts()
    
    print("Initializing PowerPoint presentation (16:9)...")
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    print("Building Slide 1: Title Slide...")
    build_slide_1(prs)
    
    print("Building Slide 2: The Problem...")
    build_slide_2(prs)
    
    print("Building Slide 3: Why The Problem Matters...")
    build_slide_3(prs)
    
    print("Building Slide 4: Existing Approach / Gaps...")
    build_slide_4(prs)
    
    print("Building Slide 5: Our Solution...")
    build_slide_5(prs)
    
    print("Building Slide 6: How It Works...")
    build_slide_6(prs)
    
    print("Building Slide 7: Key Features...")
    build_slide_7(prs)
    
    print("Building Slide 8: System Architecture...")
    build_slide_8(prs)
    
    print("Building Slide 9: AI / Intelligence & Scoring...")
    build_slide_9(prs)
    
    print("Building Slide 10: Data & Signals...")
    build_slide_10(prs, chart1, chart2)
    
    print("Building Slide 11: Impact & Performance...")
    build_slide_11(prs)
    
    print("Building Slide 12: Implementation Phases...")
    build_slide_12(prs)
    
    print("Building Slide 13: Scalability & Business Roadmap...")
    build_slide_13(prs)
    
    print("Building Slide 14: Competitive Advantage Matrix...")
    build_slide_14(prs)
    
    print("Building Slide 15: Conclusion...")
    build_slide_15(prs)
    
    print("Building Slide 16: Thank You & Disclaimer...")
    build_slide_16(prs)
    
    prs.save(PPTX_PATH)
    print(f"SUCCESS: PowerPoint presentation saved to '{PPTX_PATH}'!")

if __name__ == "__main__":
    main()
